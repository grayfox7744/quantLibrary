# -*- coding: utf-8 -*-
import utility.windutility as wu
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from datetime import date, timedelta
from pptx import Presentation

# tracking 50ETF historical volatility and implied volatility
startdate = '2008-1-1'
enddate = date.today()

# 豆油vs棕榈油
price1 = wu.wsd('Y.DCE', 'close', startdate, enddate)
price2 = wu.wsd('P.DCE', 'close', startdate, enddate)
logPrice1 = np.log(price1)
logPrice2 = np.log(price2)
spread = logPrice1 - logPrice2

plt.figure()
spread['close'].plot(label = u'豆油 vs 棕榈油')
mean = pd.rolling_mean(spread,60)
mean['close'].plot(label = u'60均线')
plt.legend()
plt.savefig('D:\\reports\\pic\\1.png')


plt.figure()
spread['close'][-120:].plot(label = u'豆油 vs 棕榈油')
mean = pd.rolling_mean(spread,20)
mean['close'][-120:].plot(label = u'20均线')
plt.legend(loc = 'upper left')
plt.savefig('D:\\reports\\pic\\2.png')


# 金 vs 银
price1 = wu.wsd('AU.SHF', 'close', startdate, enddate)
price2 = wu.wsd('AG.SHF', 'close', startdate, enddate)
logPrice1 = np.log(price1)
logPrice2 = np.log(price2)
spread = logPrice1 - logPrice2

plt.figure()
spread['close'].plot(label = u'金 vs 银')
mean = pd.rolling_mean(spread,60)
mean['close'].plot(label = u'60均线')
plt.legend(loc = 'uppper left')
plt.savefig('D:\\reports\\pic\\3.png')


'''
# create ppt file
prs = Presentation('D:\\reports\\templet\\templet.pptx')

# page 1
layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Weekly Report"
subtitle.text = "Generated on {:%Y-%m-%d}".format(date.today())

'''


