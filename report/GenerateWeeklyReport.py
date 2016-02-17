# -*- coding: utf-8 -*-
import utility.windutility as wu
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from datetime import date, timedelta
from pptx import Presentation

# tracking 50ETF historical volatility and implied volatility

enddate = date.today() - timedelta(2)
startdate = enddate - timedelta(60)

sh510050 = wu.wsd('000016.SH', 'close', startdate, enddate)
sh510050['ret'] = sh510050['close'].pct_change()
sh510050['vol_10d'] = pd.rolling_std(sh510050['ret'], window = 10) * np.sqrt(240)
sh510050['vol_20d'] = pd.rolling_std(sh510050['ret'], window = 20) * np.sqrt(240)

ivix = wu.wsd('IVIX.SH', 'close', startdate, enddate) / 100

# pic 1: price vs implied volatility
plt.figure()
sh510050['close'].plot(label = u'上证50')
ivix['close'].plot(label = 'IVIX', secondary_y = True)
plt.legend()
plt.savefig('D:\\reports\\pic\\1.png')

# pic 2: implied vol and rolling 5 day volatility
plt.figure()
sh510050['vol_10d'].plot(label = '10 days vol')
ivix['close'].plot(label = 'implied vol')
plt.legend()
plt.savefig('D:\\reports\\pic\\2.png')




# create ppt file
prs = Presentation('D:\\reports\\templet\\templet.pptx')

# page 1
layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Weekly Report"
subtitle.text = "Generated on {:%Y-%m-%d}".format(date.today())

# page 2
layout1 = prs.slide_layouts[1]
slide = prs.slides.add_slide(layout1)
title = slide.shapes.title
title.text = u"上证50 & IVIX "
placeholder = slide.placeholders[1]
pic = placeholder.insert_picture('D:\\reports\\pic\\1.png')

# page 3
slide = prs.slides.add_slide(layout1)
title = slide.shapes.title
title.text = u"IVIX & 10日滚动波动率"
placeholder = slide.placeholders[1]
pic = placeholder.insert_picture('D:\\reports\\pic\\2.png')

prs.save('D:\\reports\\ppt\\sample.pptx')
